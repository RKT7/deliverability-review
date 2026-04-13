from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
from datetime import date
from excel_parser import DIM_ORDER

# ── Colors ──
PURPLE = RGBColor(0x7B, 0x2D, 0x8E)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF0, 0xF7)
LIGHT_BORDER = RGBColor(0xD0, 0xC4, 0xE0)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_MED = RGBColor(0x55, 0x55, 0x55)
TEXT_LIGHT = RGBColor(0x66, 0x66, 0x66)

RAG_COLORS = {
    'GREEN': RGBColor(0x05, 0x96, 0x69),
    'AMBER': RGBColor(0xD9, 0x77, 0x06),
    'RED': RGBColor(0xDC, 0x26, 0x26)
}

RAG_BG = {
    'GREEN': RGBColor(0xEC, 0xFD, 0xF5),
    'AMBER': RGBColor(0xFF, 0xFB, 0xF0),
    'RED': RGBColor(0xFE, 0xF2, 0xF2)
}

RAG_EMOJI = {'GREEN': '🟢', 'AMBER': '🟡', 'RED': '🔴'}

DECISIONS = {
    'GREEN': 'Ready to Proceed',
    'AMBER': 'Proceed with Caution',
    'RED': 'Do Not Proceed — Escalation Required'
}


def _add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a shape to the slide."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=10, color=TEXT_DARK,
                   bold=False, italic=False, alignment=PP_ALIGN.LEFT, font_name='Calibri',
                   vertical_anchor=MSO_ANCHOR.TOP):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    # Set vertical anchor
    try:
        tf.paragraphs[0].alignment = alignment
    except:
        pass
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    
    return txBox


def _add_multi_text(slide, left, top, width, height, lines, font_size=9.5, color=TEXT_DARK,
                     font_name='Calibri', bullet=False, bold_first=False):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        prefix = '• ' if bullet else ''
        p.text = prefix + line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(2)
    
    return txBox


def build_ppt(ai_result, rags, risks):
    """Build the complete 5-slide PowerPoint deck.
    
    Args:
        ai_result: dict with AI-generated content
        rags: dict with dimRags and overall
        risks: list of risk dicts
        
    Returns:
        bytes: The .pptx file as bytes
    """
    prs = Presentation()
    
    # Set widescreen 16:9 layout
    prs.slide_width = Inches(14.33)
    prs.slide_height = Inches(8.06)
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank
    
    # ═══════════════════════════════════════════
    # SLIDE 1: Cover
    # ═══════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    
    # Purple background
    bg = sl.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE
    
    _add_text_box(sl, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
                  'Cybersecurity Strategy, Risk, and Architecture',
                  font_size=22, color=WHITE, bold=True)
    
    _add_text_box(sl, Inches(0.6), Inches(2.0), Inches(12), Inches(0.9),
                  'Deal Deliverability Review',
                  font_size=36, color=WHITE, bold=True)
    
    _add_text_box(sl, Inches(0.6), Inches(4.2), Inches(5), Inches(0.4),
                  'Prepared by: Accenture Security Practice',
                  font_size=12, color=WHITE)
    
    _add_text_box(sl, Inches(0.6), Inches(5.2), Inches(5), Inches(0.4),
                  f'Date: {date.today().strftime("%d/%m/%Y")}',
                  font_size=12, color=WHITE)
    
    _add_text_box(sl, Inches(0.6), Inches(7.3), Inches(10), Inches(0.3),
                  'Confidential | Accenture Internal Use Only',
                  font_size=9, color=RGBColor(0xCC, 0xBB, 0xDD))
    
    # ═══════════════════════════════════════════
    # SLIDE 2: Executive Summary
    # ═══════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    
    _add_text_box(sl, Inches(0.6), Inches(0.3), Inches(5), Inches(0.25),
                  'EXECUTIVE SUMMARY', font_size=10, color=DARK, bold=True)
    _add_text_box(sl, Inches(0.6), Inches(0.6), Inches(10), Inches(0.5),
                  'Deal Deliverability Review', font_size=24, color=PURPLE, bold=True)
    
    # Deal Snapshot (purple box)
    shape = _add_shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE,
                       Inches(0.5), Inches(1.4), Inches(5.2), Inches(1.8),
                       fill_color=PURPLE)
    _add_text_box(sl, Inches(0.7), Inches(1.5), Inches(3), Inches(0.25),
                  'DEAL SNAPSHOT', font_size=10, color=WHITE, bold=True)
    _add_text_box(sl, Inches(0.7), Inches(1.9), Inches(4.8), Inches(1.0),
                  f"Opportunity Value: {ai_result.get('opportunity_value', '')}",
                  font_size=10.5, color=WHITE)
    
    # RAG Status box
    overall = rags['overall']
    _add_shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE,
               Inches(6.0), Inches(1.4), Inches(7.8), Inches(2.6),
               fill_color=LIGHT_BG, line_color=LIGHT_BORDER)
    _add_text_box(sl, Inches(6.2), Inches(1.5), Inches(4), Inches(0.3),
                  'Overall RAG Status', font_size=12, color=DARK, bold=True)
    
    # RAG oval
    oval = _add_shape(sl, MSO_SHAPE.OVAL,
                      Inches(6.3), Inches(2.0), Inches(0.35), Inches(0.35),
                      fill_color=RAG_COLORS[overall])
    _add_text_box(sl, Inches(6.8), Inches(2.0), Inches(2), Inches(0.35),
                  overall, font_size=16, color=RAG_COLORS[overall], bold=True)
    _add_text_box(sl, Inches(6.3), Inches(2.45), Inches(7), Inches(0.3),
                  DECISIONS[overall], font_size=10, color=TEXT_DARK)
    _add_text_box(sl, Inches(6.3), Inches(2.85), Inches(7.2), Inches(0.9),
                  f"Key Justification: {ai_result.get('key_justification', '')}",
                  font_size=9.5, color=DARK, bold=True)
    
    # Deal Overview (purple box)
    _add_shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE,
               Inches(0.5), Inches(3.5), Inches(5.2), Inches(3.3),
               fill_color=PURPLE)
    _add_text_box(sl, Inches(0.7), Inches(3.6), Inches(3), Inches(0.25),
                  'DEAL OVERVIEW', font_size=10, color=WHITE, bold=True)
    overview_text = '\n'.join(ai_result.get('deal_overview', []))
    _add_text_box(sl, Inches(0.7), Inches(4.0), Inches(4.8), Inches(2.6),
                  overview_text, font_size=9.5, color=WHITE)
    
    # Positive Notes
    _add_shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE,
               Inches(6.0), Inches(4.3), Inches(7.8), Inches(2.5),
               fill_color=LIGHT_BG, line_color=LIGHT_BORDER)
    _add_text_box(sl, Inches(6.2), Inches(4.4), Inches(4), Inches(0.25),
                  'Positive Notes:', font_size=10, color=DARK, bold=True)
    _add_multi_text(sl, Inches(6.3), Inches(4.8), Inches(7.2), Inches(1.8),
                    ai_result.get('positive_notes', []), bullet=True, color=DARK)
    
    # Legend
    _add_text_box(sl, Inches(0.5), Inches(7.2), Inches(13), Inches(0.3),
                  '🟢 GREEN: Ready to Proceed     🟡 AMBER: Proceed with Caution     🔴 RED: Do Not Proceed',
                  font_size=8, color=TEXT_LIGHT)
    
    # ═══════════════════════════════════════════
    # SLIDE 3: Deliverability Assessment
    # ═══════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    
    _add_text_box(sl, Inches(0.6), Inches(0.25), Inches(6), Inches(0.2),
                  'DELIVERABILITY ASSESSMENT', font_size=10, color=DARK, bold=True)
    _add_text_box(sl, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
                  'Key Deliverability Dimensions', font_size=20, color=PURPLE, bold=True)
    _add_text_box(sl, Inches(0.6), Inches(0.95), Inches(13), Inches(0.25),
                  'Assessment across five critical dimensions. Each is independently rated GREEN / AMBER / RED.',
                  font_size=8.5, color=TEXT_MED)
    
    # Dimension card positions (3 top, 2 bottom)
    positions = [
        (0.3, 1.35, 4.4, 2.35),
        (4.9, 1.35, 4.4, 2.35),
        (9.5, 1.35, 4.4, 2.35),
        (2.5, 3.95, 4.4, 2.2),
        (7.1, 3.95, 4.4, 2.2)
    ]
    
    dimensions = ai_result.get('dimensions', [])
    for i, (x, y, w, h) in enumerate(positions):
        if i >= len(dimensions):
            break
        
        dim = dimensions[i]
        dim_name = DIM_ORDER[i] if i < len(DIM_ORDER) else dim.get('name', '')
        dim_rag = rags['dimRags'].get(dim_name, 'GREEN')
        
        # Card outline
        _add_shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE,
                   Inches(x), Inches(y), Inches(w), Inches(h),
                   fill_color=WHITE, line_color=PURPLE, line_width=1.5)
        
        # Dimension title
        _add_text_box(sl, Inches(x + 0.15), Inches(y + 0.08), Inches(w - 0.3), Inches(0.35),
                      f'Dimension {i + 1}: {dim_name}',
                      font_size=8.5, color=PURPLE, bold=True)
        
        # Bullets
        bullets = dim.get('bullets', [])
        _add_multi_text(sl, Inches(x + 0.15), Inches(y + 0.5), Inches(w - 0.3), Inches(h - 1.15),
                        bullets, font_size=7.5, color=TEXT_DARK, bullet=True)
        
        # RAG + Comments footer
        _add_text_box(sl, Inches(x + 0.1), Inches(y + h - 0.4), Inches(w - 0.2), Inches(0.3),
                      f"RAG: {dim_rag}  |  {dim.get('comments', '')}",
                      font_size=7, color=DARK, bold=True)
    
    # RED summary box
    _add_shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE,
               Inches(0.3), Inches(6.5), Inches(6.6), Inches(1.1),
               fill_color=RAG_BG['RED'], line_color=RAG_COLORS['RED'])
    _add_text_box(sl, Inches(0.5), Inches(6.6), Inches(5), Inches(0.2),
                  'Mandatory Gates / Critical Blockers',
                  font_size=9, color=RAG_COLORS['RED'], bold=True)
    _add_text_box(sl, Inches(0.5), Inches(6.85), Inches(6.2), Inches(0.6),
                  ai_result.get('red_summary', ''), font_size=8.5, color=TEXT_DARK)
    
    # AMBER summary box
    _add_shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE,
               Inches(7.1), Inches(6.5), Inches(6.6), Inches(1.1),
               fill_color=RAG_BG['AMBER'], line_color=RAG_COLORS['AMBER'])
    _add_text_box(sl, Inches(7.3), Inches(6.6), Inches(5), Inches(0.2),
                  'Required Actions for AMBER Items',
                  font_size=9, color=RAG_COLORS['AMBER'], bold=True)
    _add_text_box(sl, Inches(7.3), Inches(6.85), Inches(6.2), Inches(0.6),
                  ai_result.get('amber_summary', ''), font_size=8.5, color=TEXT_DARK)
    
    # ═══════════════════════════════════════════
    # SLIDE 4: Risks, Assumptions, Next Steps
    # ═══════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    
    _add_text_box(sl, Inches(0.6), Inches(0.25), Inches(6), Inches(0.2),
                  'CRITICAL RISKS & WAY FORWARD', font_size=10, color=DARK, bold=True)
    _add_text_box(sl, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
                  'Critical Risks, Key Assumptions & Way Forward',
                  font_size=20, color=PURPLE, bold=True)
    
    # Risk cards (up to 3)
    for i, risk in enumerate(risks[:3]):
        ry = 1.3 + i * 1.6
        _add_shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE,
                   Inches(0.3), Inches(ry), Inches(6.5), Inches(1.4),
                   fill_color=WHITE, line_color=PURPLE, line_width=1.5)
        _add_text_box(sl, Inches(0.5), Inches(ry + 0.1), Inches(3), Inches(0.2),
                      f'Risk {i + 1}', font_size=10, color=DARK, bold=True)
        _add_text_box(sl, Inches(0.5), Inches(ry + 0.4), Inches(6), Inches(0.3),
                      risk.get('risk', ''), font_size=8.5, color=TEXT_DARK)
        _add_text_box(sl, Inches(0.5), Inches(ry + 0.8), Inches(6), Inches(0.45),
                      f"Mitigation: {risk.get('mit', 'N/A')}",
                      font_size=8.5, color=TEXT_MED, italic=True)
    
    # Assumptions box
    _add_shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE,
               Inches(7.2), Inches(1.3), Inches(6.8), Inches(3.2),
               fill_color=LIGHT_BG, line_color=LIGHT_BORDER)
    _add_text_box(sl, Inches(7.4), Inches(1.4), Inches(3), Inches(0.25),
                  'KEY ASSUMPTIONS', font_size=10, color=PURPLE, bold=True)
    
    assumptions = ai_result.get('assumptions', [])
    for i, assumption in enumerate(assumptions[:4]):
        _add_text_box(sl, Inches(7.5), Inches(1.8 + i * 0.7), Inches(6.3), Inches(0.6),
                      f"{i + 1}. {assumption}", font_size=8.5, color=TEXT_DARK)
    
    # Next Steps box
    _add_shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE,
               Inches(7.2), Inches(4.8), Inches(6.8), Inches(2.5),
               fill_color=LIGHT_BG, line_color=LIGHT_BORDER)
    _add_text_box(sl, Inches(7.4), Inches(4.9), Inches(4), Inches(0.25),
                  'NEXT STEPS & ACTION PLAN', font_size=10, color=PURPLE, bold=True)
    
    next_steps = ai_result.get('next_steps', [])
    for i, ns in enumerate(next_steps[:2]):
        _add_text_box(sl, Inches(7.5), Inches(5.3 + i * 1.0), Inches(6.3), Inches(0.8),
                      f"{ns['title']} — {ns['desc']}\nOwner: {ns['owner']}",
                      font_size=8.5, color=DARK)
    
    # ═══════════════════════════════════════════
    # SLIDE 5: Thank You
    # ═══════════════════════════════════════════
    sl = prs.slides.add_slide(blank_layout)
    
    bg = sl.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PURPLE
    
    _add_text_box(sl, Inches(1), Inches(2.5), Inches(12), Inches(1.5),
                  'THANK YOU', font_size=48, color=WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)
    _add_text_box(sl, Inches(1), Inches(5), Inches(12), Inches(0.5),
                  'Confidential | Accenture Internal Use Only | © 2026 Accenture',
                  font_size=12, color=RGBColor(0xCC, 0xBB, 0xDD),
                  alignment=PP_ALIGN.CENTER)
    
    # ── Save to bytes ──
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
